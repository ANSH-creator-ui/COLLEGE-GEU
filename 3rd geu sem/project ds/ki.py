from pptx import Presentation
from pptx.util import Inches, Pt

# Create a presentation
prs = Presentation()

# --- Slide 1: Title Slide ---
slide1 = prs.slides.add_slide(prs.slide_layouts[0])
title = slide1.shapes.title
subtitle = slide1.placeholders[1]

title.text = "Crystallization Method of Making a Hand Warmer"
subtitle.text = "A Science Project Presentation\nBy: [Your Name]"


# --- Slide 2: Introduction ---
slide2 = prs.slides.add_slide(prs.slide_layouts[1])
slide2.shapes.title.text = "Introduction"
content = slide2.placeholders[1]
content.text = (
    "Hand warmers use the process of crystallization to produce heat.\n"
    "A supersaturated solution of sodium acetate releases energy when it crystallizes.\n"
    "This exothermic process provides instant warmth."
)

# --- Slide 3: Materials Required ---
slide3 = prs.slides.add_slide(prs.slide_layouts[1])
slide3.shapes.title.text = "Materials Required"
content = slide3.placeholders[1]
content.text = (
    "• Sodium acetate trihydrate\n"
    "• Water\n"
    "• Heat-safe beaker\n"
    "• Metal spoon or stirrer\n"
    "• Plastic pouch or sealed container\n"
    "• Small metal disk (for triggering crystallization)"
)

# --- Slide 4: Procedure ---
slide4 = prs.slides.add_slide(prs.slide_layouts[1])
slide4.shapes.title.text = "Procedure"
content = slide4.placeholders[1]
content.text = (
    "1. Heat water and dissolve sodium acetate until no more can dissolve (saturated solution).\n"
    "2. Allow the solution to cool without disturbing — it becomes supersaturated.\n"
    "3. Pour the solution into a sealed plastic pouch with a small metal disk inside.\n"
    "4. To activate, flex or click the metal disk — crystallization begins, releasing heat.\n"
    "5. To reuse, place the pouch in boiling water until crystals dissolve again."
)

# --- Slide 5: Working Principle ---
slide5 = prs.slides.add_slide(prs.slide_layouts[1])
slide5.shapes.title.text = "Working Principle"
content = slide5.placeholders[1]
content.text = (
    "• The sodium acetate solution is supersaturated — unstable but appears liquid.\n"
    "• Clicking the disk provides nucleation points for crystals to form.\n"
    "• Crystallization releases stored heat energy (exothermic reaction).\n"
    "• The same process can be reversed by reheating to dissolve the crystals."
)

# --- Slide 6: Advantages and Applications ---
slide6 = prs.slides.add_slide(prs.slide_layouts[1])
slide6.shapes.title.text = "Advantages and Applications"
content = slide6.placeholders[1]
content.text = (
    "• Reusable and environmentally friendly.\n"
    "• Safe for human use, no combustion involved.\n"
    "• Used in camping, medical therapy, and cold environments.\n"
    "• Demonstrates practical use of crystallization and exothermic reactions."
)

# --- Slide 7: Conclusion ---
slide7 = prs.slides.add_slide(prs.slide_layouts[1])
slide7.shapes.title.text = "Conclusion"
content = slide7.placeholders[1]
content.text = (
    "The crystallization method of hand warmers is a brilliant application of chemistry.\n"
    "It shows how energy stored in a chemical solution can be released safely and reused.\n"
    "This simple yet effective device converts theory into everyday usefulness."
)

# Save the PowerPoint file
prs.save("Crystallization_Hand_Warmer_Presentation.pptx")
print("✅ Presentation created successfully: Crystallization_Hand_Warmer_Presentation.pptx")